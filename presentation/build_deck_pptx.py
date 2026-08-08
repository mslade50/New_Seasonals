"""Build a native, editable PowerPoint (.pptx) from deck_spec.json.

Text (titles, bullets, callouts) is native and editable; speaker notes go into the
notes pane; each figure is embedded as its cropped 2x image (figs/<id>.png) so the
slides match the HTML deck. Run build_figures.py + screenshot + crop_figs.py first.

Usage:
    python presentation/build_deck_pptx.py [spec.json] [out.pptx]
    # no args -> deck_spec.json -> The_Book.pptx (the full deck)
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
SPEC_PATH = Path(sys.argv[1]) if len(sys.argv) > 1 else HERE / "deck_spec.json"
SPEC = json.loads(SPEC_PATH.read_text(encoding="utf-8"))
FIGDIR = HERE / "figs"
OUT = Path(sys.argv[2]) if len(sys.argv) > 2 else HERE / "The_Book.pptx"

# Palette (matches the HTML deck)
SLIDE_BG = RGBColor(0x0C, 0x16, 0x26)
INK = RGBColor(0xF4, 0xF9, 0xFF)
SOFT = RGBColor(0xC2, 0xD0, 0xE2)
MUTED = RGBColor(0x8A, 0xA0, 0xB8)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
CALLOUT_BG = RGBColor(0x10, 0x24, 0x2A)
CALLOUT_TX = RGBColor(0xD6, 0xFF, 0xF8)
PILL_TX = RGBColor(0x04, 0x22, 0x1F)
FONT = "Segoe UI"

EMU = 914400
SW, SH = 13.333, 7.5
# Right region for figures
RX, RY, RW, RH = 5.85, 2.45, 7.0, 4.55


def _font(run, size, color, bold=False, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_bg(slide):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    r.fill.solid()
    r.fill.fore_color.rgb = SLIDE_BG
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tf


def para(tf, first=False):
    return tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()


def add_figure(slide, fig_id):
    img = FIGDIR / f"{fig_id}.png"
    if not img.exists():
        return
    iw, ih = Image.open(img).size
    ar = iw / ih
    w = RW
    h = w / ar
    if h > RH:
        h = RH
        w = h * ar
    left = RX + (RW - w) / 2
    top = RY + (RH - h) / 2
    slide.shapes.add_picture(str(img), Inches(left), Inches(top), Inches(w), Inches(h))


def add_notes(slide, text):
    if not text:
        return
    slide.notes_slide.notes_text_frame.text = text


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    tf = textbox(slide, 0.9, 1.25, 11, 0.4)
    p = para(tf, True)
    r = p.add_run(); r.text = "QUANTITATIVE EQUITY TRADING SYSTEM"
    _font(r, 13, TEAL, bold=True)

    tf = textbox(slide, 0.88, 1.7, 11.5, 1.5)
    p = para(tf, True)
    r = p.add_run(); r.text = SPEC["deck_title"].split(":")[0]
    _font(r, 60, INK, bold=True)

    tf = textbox(slide, 0.9, 3.45, 9.5, 1.0)
    p = para(tf, True)
    r = p.add_run(); r.text = SPEC["subtitle"]
    _font(r, 22, SOFT)

    if SPEC.get("tagline"):
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.75), Inches(7.1), Inches(0.66))
        pill.fill.solid(); pill.fill.fore_color.rgb = TEAL
        pill.line.fill.background(); pill.shadow.inherit = False
        tf = pill.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = SPEC.get("tagline", "")
        _font(r, 15, PILL_TX, bold=True)

    tf = textbox(slide, 0.9, 5.75, 11.0, 1.4)
    p = para(tf, True); p.line_spacing = 1.35
    r = p.add_run(); r.text = SPEC.get("narrative_arc", "")
    _font(r, 12.5, MUTED)


def content(prs, s, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    tf = textbox(slide, 0.62, 0.5, 11, 0.34)
    p = para(tf, True)
    r = p.add_run(); r.text = s["section"].upper()
    _font(r, 12, TEAL, bold=True)

    tf = textbox(slide, 0.6, 0.82, 12.2, 1.0)
    p = para(tf, True); p.line_spacing = 0.98
    r = p.add_run(); r.text = s["title"]
    _font(r, 31, INK, bold=True)

    if s.get("subtitle"):
        tf = textbox(slide, 0.62, 1.82, 11.5, 0.5)
        p = para(tf, True)
        r = p.add_run(); r.text = s["subtitle"]
        _font(r, 14, MUTED)

    # Bullets
    bullets = s.get("bullets") or []
    if bullets:
        tf = textbox(slide, 0.62, 2.5, 4.95, 3.1)
        for i, b in enumerate(bullets):
            p = para(tf, first=(i == 0))
            p.space_after = Pt(8); p.line_spacing = 1.04
            rb = p.add_run(); rb.text = "▪  "; _font(rb, 12, TEAL, bold=True)
            rt = p.add_run(); rt.text = b; _font(rt, 13, SOFT)

    # Callout
    if s.get("callout"):
        co = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(5.75), Inches(4.95), Inches(1.05))
        co.fill.solid(); co.fill.fore_color.rgb = CALLOUT_BG
        co.line.color.rgb = TEAL; co.line.width = Pt(1.0)
        co.shadow.inherit = False
        tf = co.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Pt(12)
        p = tf.paragraphs[0]; p.line_spacing = 1.1
        r = p.add_run(); r.text = s["callout"]
        _font(r, 15, CALLOUT_TX, bold=True)

    add_figure(slide, s.get("figure", ""))

    # Footer slide number
    tf = textbox(slide, 11.6, 7.05, 1.5, 0.3)
    p = para(tf, True); p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f"{idx} / {total}"
    _font(r, 10, MUTED)

    add_notes(slide, s.get("speaker_notes", ""))


def main():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    cover(prs)
    items = [s for s in SPEC["slides"] if s.get("section") != "Open"]
    for i, s in enumerate(items, start=1):
        content(prs, s, i, len(items))
    prs.save(OUT)
    print(f"Wrote {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
